import streamlit as st
from pptx import Presentation
import io
import zipfile

st.set_page_config(page_title="Bricks 4 Kidz 证书助手", layout="centered")
st.title("🎓 证书批量生成器")

uploaded_file = st.file_uploader("1. 上传 PPT 模板", type="pptx")
new_date = st.text_input("2. 输入新日期", "2026年3月1日")
names_text = st.text_area("3. 输入名字名单（每行一个）", "小明\n小红")

if uploaded_file and st.button("开始生成"):
    name_list = [n.strip() for n in names_text.split("\n") if n.strip()]
    zip_buffer = io.BytesIO()
    
    with zipfile.ZipFile(zip_buffer, "a", zipfile.ZIP_DEFLATED) as zip_file:
        for name in name_list:
            prs = Presentation(io.BytesIO(uploaded_file.getvalue()))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        # 查找并替换日期
                        if "Date" in shape.text:
                            shape.text = f"Date: {new_date}"
                        # 查找并替换姓名（根据高度锁定中下部区域）
                        if 5.0 < shape.top.inches < 7.0:
                            if "Presented" not in shape.text:
                                shape.text = name
            
            ppt_io = io.BytesIO()
            prs.save(ppt_io)
            zip_file.writestr(f"证书_{name}.pptx", ppt_io.getvalue())

    st.success("全部生成成功！")
    st.download_button("点击下载 ZIP 包", data=zip_buffer.getvalue(), file_name="certs.zip")
